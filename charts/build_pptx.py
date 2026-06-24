#!/usr/bin/env python3
"""Sestaví prezentaci ZJISTENI z decku + grafů (charts/*.png).
Spuštění:  python3 charts/build_pptx.py  → Respekt_zjisteni.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CH = os.path.join(ROOT, "charts")
OUT = os.path.join(ROOT, "Respekt_zjisteni.pptx")

DARK = RGBColor(0x1F, 0x2A, 0x37)
RED = RGBColor(0xC0, 0x39, 0x2B)
INK = RGBColor(0x22, 0x28, 0x31)
MUTED = RGBColor(0x6B, 0x72, 0x80)
TINT = RGBColor(0xF4, 0xF5, 0xF7)
REDTINT = RGBColor(0xF8, 0xEC, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xD7, 0xDB, 0xE0)
BODY = "Calibri"
HEAD = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def txt(s, text, x, y, w, h, size=14, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, font=BODY, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text if isinstance(text, list) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = ln
        f = r.font
        f.size = Pt(size)
        f.name = font
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def bullets(s, items, x, y, w, h, size=14, color=INK, gap=6):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        rb = p.add_run(); rb.text = "•  "
        rb.font.size = Pt(size); rb.font.name = BODY; rb.font.color.rgb = RED; rb.font.bold = True
        rt = p.add_run(); rt.text = it
        rt.font.size = Pt(size); rt.font.name = BODY; rt.font.color.rgb = color
    return tb


def card(s, x, y, w, h, fill, shadow=True):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = 0.04
    except Exception:
        pass
    return sp


def img_fit(s, path, bx, by, bw, bh):
    with Image.open(path) as im:
        iw, ih = im.size
    r = iw / ih
    w, h = bw, bw / r
    if h > bh:
        h, w = bh, bh * r
    x = bx + (bw - w) / 2
    y = by + (bh - h) / 2
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))


def cotim(s, text, x=0.6, y=5.35, w=12.13, h=1.55):
    card(s, x, y, w, h, REDTINT)
    txt(s, "CO S TÍM", x + 0.28, y + 0.16, 3, 0.3, size=11, color=RED, bold=True)
    txt(s, text, x + 0.28, y + 0.5, w - 0.56, h - 0.6, size=14.5, color=INK)


def kicker(s, n):
    txt(s, f"ZJIŠTĚNÍ {n:02d}", 0.62, 0.34, 4, 0.32, size=12, color=RED, bold=True)


# ---------- obsah ----------
F = [
    dict(n=1, h="Báze je spokojená a loajální — kotvou je hodnota, ne cena",
         ev=["92 % udává vysokou/velmi vysokou pravděpodobnost setrvání; jen 12,8 % zvažovalo odchod.",
             "Spokojenost s atributy je 90–99 % napříč.",
             "Chvála stojí na šíři témat (461), hloubce (385) a kvalitě psaní (370)."],
         img="chvala.png",
         co="Chránit jádro (hloubka, šíře, nezávislost, kvalita psaní) a komunikovat hodnotu, ne slevu."),
    dict(n=2, h="Retence se láme v prvním roce",
         ev=["Churn < 1 rok 17 % a 1–5 let 16 % vs. 16–20 let jen 5 %.",
             "Noví předplatitelé (151) jsou nejkřehčí fáze."],
         img="churn_tenure.png",
         co="Onboarding prvního roku — dovést k návyku (archiv, audio, hloubka). Měřit retenci v prvních 12 měsících."),
    dict(n=3, h="Zapojení = retence; pasivní digitál je tichý odchod",
         ev=["Pasivní digitál (178): churn 19 % vs. digitální jádro (688): 10 %.",
             "Web vůbec nepoužívá 23–27 % mladších kohort."],
         img="churn_segments.png",
         co="Frekvenci návštěv brát jako early-warning churnu; reaktivovat utlumené (newsletter, notifikace, audio)."),
    dict(n=4, h="Mladší muži jsou nejrizikovější kohorta — táhne to produkt",
         ev=["Churn mladší muži 18 % vs. starší ženy 8,5 %.",
             "Nejvíc je trápí audio (21 %) a UX / ovládání (16 %).",
             "„Nic nevadí“ řekne jen 5 % mladších mužů — nejkritičtější."],
         img="heat_riziko.png",
         co="Udržení mladších = produktové investice (audio, aplikace, UX), ne další obsah."),
    dict(n=5, h="Audio je jediná slabina spokojenosti — a největší příležitost",
         ev=["Audio nejnižší spokojenost: 83 % (mladší muži 72 %) vs. 95–99 % u ostatních atributů.",
             "Výtky k audiu / AI hlasu 102; volný text přidal +20 zmínek.",
             "1 026 lidí audio nevyužívá; dalších 119 poslech v appce nezkusilo."],
         img="heat_spokojenost.png",
         co="Zlepšit kvalitu AI hlasu (výslovnost, rozlišení hlasů) a aktivovat ty, kdo audio zatím míjejí."),
    dict(n=6, h="Vyhledávání a archiv = nejkonkrétnější funkční mezera",
         ev=["Vyhledávání: výtka 32×, nevyžádaně v bariérách 13×, přání v appce 229.",
             "Přehlednost / archiv / orientace: 36 výtek."],
         img="vyhledavani.png",
         co="Vyhledávání a orientaci v archivu vysoko na app/web roadmapu — opakovaná a řešitelná rychlá výhra."),
    dict(n=7, h="Tisk není setrvačnost — je to rituál a kotva retence",
         ev=["Část digitál aktivně odmítá: preferuje tisk (15), digitál jako doplněk (20).",
             "„Tištěné posílám rodině“ (12) — opakovaný vzorec."],
         img="tisk_ritual.png",
         co="Nehnat všechny do digitálu; chránit tištěný zážitek. Potenciál pro dárkové a rodinné předplatné."),
    dict(n=8, h="Obsahové výtky míří na kulturu, jednostrannost a tón",
         ev=["Kulturní rubrika 85, jednostrannost / bias 64, délka 67, tón 28.",
             "Hodnotové souznění je v chvále (173) i ve výtce (bias 64) — stejná osa, opačná valence."],
         img="vytky.png",
         co="Redakční reflexe kultury, vyváženosti a délky; „bias“ je menšinový, ale hlasitý a hodnotově nabitý signál."),
    dict(n=9, h="Mladé publikum chce jiný obsah; akvizici táhne mise + sleva",
         ev=["Mladší ženy: reportáže 42 %, rodina a vztahy 15 %, míň politiky; mladší muži politika 44 %.",
             "Konverze mladých: podpora médií 73 % a sleva 22 % (vs. 9–12 % u starších)."],
         img="heat_zajmy.png",
         co="Cílený obsah a marketing pro mladší; slevová akvizice funguje, ale bez onboardingu prvního roku je ztratíme."),
    dict(n=10, h="Konkurence je Deník N a veřejnoprávní; publikum je náročné",
         ev=["Nejčastější jiný zdroj: Deník N 549, iRozhlas 414, Seznam 403, ČT 374.",
             "Předplatitelé kombinují domácí + veřejnoprávní + zahraniční zdroje."],
         img="zdroje.png",
         co="Sledovat Deník N jako přímého konkurenta; odlišení stavět na hloubce, kurátorství a překladech (66)."),
    # 11–20
    dict(n=11, h="Cena není churn páka — citlivost na cenu je minimální",
         ev=["Cena / paywall ve výtkách jen 11× — prakticky na dně 16 témat.",
             "Sleva jako konverzní motiv hraje roli hlavně u mladých a nových."],
         stat="11×", statlab="zmínek o ceně ve výtkách\n(z 1 307 odpovědí)",
         co="Zdražení nese menší riziko, než se obvykle čeká, pokud zůstane hodnota. Slevu cíleně na mladé, ne plošně."),
    dict(n=12, h="Publikum chce delší obsah, ne kratší",
         ev=["U textů: 67 % spíše delší, jen 17 % kratší → 4 z 5 lidí s názorem chtějí delší.",
             "Audioverze 64 %, podcasty 61 % (z těch s názorem).",
             "Proti tomu „moc dlouhé“ jen 67 výtek = hlasitá menšina."],
         img="delka_preference.png",
         co="Nepřeklápět produkt ke krátkému obsahu; zkracovat selektivně. U videa naopak prostor pro kratší formát."),
    dict(n=13, h="Kdo odchází kvůli obsahu, vadí mu jednostrannost a tón",
         ev=["Churneři over-index: bias 6 vs. 3 %, dále délka, kultura, tón.",
             "Dvě churn-pružiny: produktová (mladší muži) a hodnotově-redakční."],
         img="churneri_overindex.png",
         co="Retence musí mířit na obě osy zvlášť — produktové opravy neudrží toho, kdo odchází kvůli vnímané jednostrannosti."),
    dict(n=14, h="Homepage není vstupní bod — distribuce přes newsletter a sítě",
         ev=["Část přichází přes externí odkaz (newsletter, sítě, RSS, QR z tisku) — 15.",
             "Část homepage prakticky nepoužívá — 8 (jde rovnou na audio, RSS, papír)."],
         stat="15 + 8", statlab="přes odkaz / homepage nepoužívá\n→ newsletter a sítě = vstupní brána",
         co="Brát newsletter a sociální sítě jako distribuční kanál, ne marketingovou ozdobu — investovat do nich."),
    dict(n=15, h="Audio a offline jsou tahouny přechodu na digitál",
         ev=["Důvody přechodu: pohodlí 529, audio jen digitálně 344, praktické 298, ekologie 200.",
             "Offline: stáhnout vydání 366, „lepší offline režim“ 93."],
         img="digital_drivers.png",
         co="Audio a offline čtení = on-ramp na digitál; propagovat je jako důvod vyzkoušet appku, doladit offline režim."),
    dict(n=16, h="Přístupnost (velikost písma) je snadná výhra vzhledem k věku",
         ev=["Nevyžádaně v poli „co chybí“: přístupnost (velikost písma, zoom, čtení bez brýlí) 10×.",
             "Věkový profil báze je starší — dopad je nadproporční."],
         stat="45+", statlab="věk většiny báze —\ndopad přístupnosti je velký",
         co="Nastavitelná velikost písma a zoom obrázků = low-effort / high-fit zásah do aplikace."),
    dict(n=17, h="Top přání v aplikaci: odlišit přečtené, offline, vlastní playlist",
         ev=["Odlišení přečteného 251, offline 366, otevírat ze sítí 126, playlist 120, CarPlay 118.",
             "Z volného textu: přehlednost UI (12), výkon / stabilita (10)."],
         img="app_prani.png",
         co="App roadmapa: orientace ve vydání (co už jsem četl), offline, audio playlist a CarPlay."),
    dict(n=18, h="Poslech v aplikaci prohrává s agregátory — cíl jsou ti, kdo nezkusili",
         ev=["Důvod není chyba: jiná platforma (Spotify/Apple) 126, zvyk 30, agregace 65.",
             "119 lidí poslech v appce vůbec nezkusilo; tech. problémy (21) jsou menšina."],
         img="poslech_app.png",
         co="Nepřetahovat ze Spotify; cílit na 119 nevyzkoušejících a nabídnout, co agregátor neumí (návaznost na text)."),
    dict(n=19, h="Loajalita stojí na dvou osách: důvěra/fakta vs. vyváženost",
         ev=["Důvěra / ověřená fakta (194) a vyváženost / objektivita (170) jsou dva různé důvody.",
             "Serióznost vs. nestrannost úhlů — rezonují u různých lidí."],
         stat="194 : 170", statlab="důvěra/fakta vs. vyváženost\n— dva různé důvody loajality",
         co="V komunikaci nezaměňovat „věříme faktům“ a „dáváme různé úhly“; vyváženost je táž osa, kterou jiní kritizují jako bias."),
    dict(n=20, h="Drobné, ale opakované značkové signály: obálka a tón",
         ev=["Grafika / obálka / ilustrace 39 — opakované „chybí Reisenauer“.",
             "Tón / víc pozitivního 28 — „depresivní“, chybí naděje a řešení."],
         stat="39 + 28", statlab="obálka/ilustrace + „depresivní“ tón\n— emoční vazba ke značce",
         co="Vizuální identita má emoční vazbu ke značce; zvážit konstruktivní / řešení-orientovaný obsah jako protiváhu."),
]


def finding_slide(d):
    s = slide()
    kicker(s, d["n"])
    txt(s, d["h"], 0.6, 0.66, 12.13, 1.0, size=25, color=INK, bold=True, font=HEAD, spacing=0.98)
    if d.get("img"):
        bullets(s, d["ev"], 0.62, 1.95, 5.55, 3.1, size=14.5)
        img_fit(s, os.path.join(CH, d["img"]), 6.45, 1.8, 6.35, 3.35)
    else:
        card(s, 0.62, 1.95, 4.7, 2.9, TINT)
        sl = len(d["stat"])
        ssize = 46 if sl <= 4 else (40 if sl <= 7 else 34)
        txt(s, d["stat"], 0.62, 2.2, 4.7, 1.5, size=ssize, color=RED, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD)
        txt(s, d["statlab"].split("\n"), 0.62, 3.62, 4.7, 1.1, size=13, color=MUTED,
            align=PP_ALIGN.CENTER)
        bullets(s, d["ev"], 5.7, 2.1, 7.1, 2.7, size=15)
    cotim(s, d["co"])


# ---------- titulní ----------
def title_slide():
    s = slide(DARK)
    txt(s, "PŘEDPLATITELSKÝ PRŮZKUM", 0.9, 1.5, 11, 0.5, size=15, color=RGBColor(0xC9, 0xB0, 0xAC), bold=True)
    txt(s, "Klíčová zjištění pro management", 0.9, 2.15, 11.5, 1.8, size=46, color=WHITE, bold=True, font=HEAD, spacing=1.0)
    txt(s, "20 zjištění ve formátu: zjištění → evidence → co s tím", 0.92, 4.0, 11, 0.5, size=18, color=RGBColor(0xC9, 0xCF, 0xD6))
    txt(s, "N = 2 139 dokončených odpovědí   ·   export červen 2026", 0.92, 6.5, 11, 0.4, size=13, color=RGBColor(0x9A, 0xA3, 0xAD))


def exec_slide():
    s = slide()
    txt(s, "Exekutivní shrnutí", 0.6, 0.45, 8, 0.8, size=30, color=INK, bold=True, font=HEAD)
    frames = [
        ("Báze je spokojená a loajální.", "92 % vysoké setrvání, jen 12,8 % zvažuje odchod. Kotvou je hodnota a kvalita, ne cena."),
        ("Riziko je v rozložení, ne v průměru.", "Churn se koncentruje: první rok (17 %), pasivní digitál (19 %), mladší muži (18 %)."),
        ("Tři produktové páky.", "Audio / AI hlas, vyhledávání a archiv, přehlednost aplikace."),
        ("Redakční signály jsou menšinové, ale hlasité.", "Kultura, vnímaná jednostrannost, „depresivní“ tón, délka."),
    ]
    y = 1.5
    for head, body in frames:
        txt(s, head, 0.62, y, 6.5, 0.4, size=16, color=RED, bold=True)
        txt(s, body, 0.62, y + 0.42, 6.5, 0.8, size=13.5, color=INK, spacing=1.0)
        y += 1.28
    # karta doporučení
    card(s, 7.4, 1.5, 5.35, 5.4, TINT)
    txt(s, "TOP 5 DOPORUČENÍ", 7.7, 1.75, 4.8, 0.4, size=14, color=RED, bold=True)
    recs = [
        "Onboarding prvního roku — dovést k návyku.",
        "Aktivace utlumených (frekvence = early-warning).",
        "Audio jako priorita — zlepšit hlas i aktivovat.",
        "Vyhledávání + přehlednost do app roadmapy.",
        "Chránit jádro a komunikovat hodnotu, ne slevu.",
    ]
    tb = s.shapes.add_textbox(Inches(7.7), Inches(2.35), Inches(4.8), Inches(4.3))
    tf = tb.text_frame; tf.word_wrap = True
    for i, r in enumerate(recs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12); p.line_spacing = 1.04
        rn = p.add_run(); rn.text = f"{i+1}.  "
        rn.font.size = Pt(15); rn.font.bold = True; rn.font.color.rgb = RED; rn.font.name = BODY
        rt = p.add_run(); rt.text = r
        rt.font.size = Pt(14.5); rt.font.color.rgb = INK; rt.font.name = BODY


def divider_slide(title, sub):
    s = slide(DARK)
    txt(s, sub, 0.9, 2.7, 11, 0.5, size=16, color=RGBColor(0xC9, 0xB0, 0xAC), bold=True)
    txt(s, title, 0.9, 3.2, 11.5, 1.2, size=40, color=WHITE, bold=True, font=HEAD)


def closing_slide():
    s = slide(DARK)
    txt(s, "Tři priority do akce", 0.9, 1.2, 11, 1.0, size=38, color=WHITE, bold=True, font=HEAD)
    items = [
        ("Udržet nové a utlumené", "Onboarding prvního roku + aktivace pasivního digitálu. Frekvence návštěv jako early-warning."),
        ("Produktové páky", "Audio / AI hlas, vyhledávání a archiv, přehlednost aplikace — kde se láme udržení mladších."),
        ("Chránit jádro a hodnotu", "Hloubka, nezávislost, tištěný rituál. Komunikovat hodnotu, ne slevu; cena není churn páka."),
    ]
    y = 2.7
    for i, (head, body) in enumerate(items):
        txt(s, f"{i+1}", 0.95, y, 0.7, 0.8, size=34, color=RED, bold=True, font=HEAD)
        txt(s, head, 1.75, y, 10.5, 0.5, size=20, color=WHITE, bold=True)
        txt(s, body, 1.75, y + 0.5, 10.6, 0.7, size=14, color=RGBColor(0xC9, 0xCF, 0xD6), spacing=1.0)
        y += 1.4


# ---------- sestavení ----------
title_slide()
exec_slide()
for d in F[:10]:
    finding_slide(d)
divider_slide("Hlubší zjištění", "11–20  ·  JEMNĚJŠÍ VZORCE A PRODUKTOVÉ DETAILY")
for d in F[10:]:
    finding_slide(d)
closing_slide()

prs.save(OUT)
print(f"Uloženo: {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} snímků)")
